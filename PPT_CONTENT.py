from pptx import Presentation

def split_content_with_threshold(text, threshold):
    pullstops = text.split('.')
    result = []
    chunk = ''
    for i, pullstop in enumerate(pullstops):
        chunk += pullstop.strip() + '.'
        if (i + 1) % threshold == 0:
            result.append(chunk.strip())
            chunk = ''
    if chunk:
        result.append(chunk.strip())
    return result

def create_ppt_from_chunks(chunks):
    prs = Presentation()
    for chunk in chunks:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = 'Chunk'
        slide.placeholders[1].text = chunk
    return prs

text = """
C Programming is a foundational language in the world of computer science and software development. It was developed in the early 1970s by Dennis Ritchie at Bell Labs and has since become one of the most widely used and influential programming languages. Known for its efficiency, flexibility, and portability, C has been the basis for the development of many other languages, including C++, Java, and Python. Its syntax is relatively simple and straightforward, making it an excellent choice for beginners to learn programming concepts. C is often used in system programming, embedded programming, and developing operating systems due to its ability to directly interact with hardware and low-level system components.

One of the key features of C Programming is its powerful support for pointers, which allow developers to directly manipulate memory addresses. While this can be challenging for novice programmers, it offers significant advantages in terms of performance and flexibility. Additionally, C provides a rich set of built-in functions and libraries that facilitate various tasks such as input/output operations, string manipulation, and mathematical computations. Its modular structure encourages the development of reusable code components, promoting code maintainability and scalability. Despite its age, C remains a vital language in the software industry, with applications ranging from embedded systems in consumer electronics to high-performance computing in scientific research. Mastering C Programming not only equips individuals with valuable skills but also fosters a deeper understanding of computer architecture and software development principles.
"""

threshold = 2
result = split_content_with_threshold(text, threshold)

prs = create_ppt_from_chunks(result)
prs.save(r'C:\Practical Freaks\content_chunks.pptx')
print("Presentation created successfully.")
